import streamlit as st
from pptx import Presentation
import io

st.title("PPTX Oluşturucu")

title = st.text_input("Slayt Başlığı")
content = st.text_area("İçerik (her satır bir madde)")

if st.button("PPTX Oluştur"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for line in content.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.level = 1

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    st.download_button(
        "PPTX Dosyasını İndir",
        buffer,
        file_name="sunum.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
